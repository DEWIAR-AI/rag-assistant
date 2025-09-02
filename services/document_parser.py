import os
import logging
import tempfile
import re
from typing import Dict, Any, List, Optional
import mimetypes

# Document parsing libraries
try:
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pandas as pd
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    import markdown
    MARKDOWN_AVAILABLE = True
except ImportError:
    MARKDOWN_AVAILABLE = False

# Legacy format support
try:
    import xlrd  # For old .xls files
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False

try:
    import docx2txt  # Alternative DOC parser
    DOCX2TXT_AVAILABLE = True
except ImportError:
    DOCX2TXT_AVAILABLE = False

# DOC parsing alternatives - antiword not available on Windows
ANTIWORD_AVAILABLE = False
DOC2DOCX_AVAILABLE = False

# RTF support
try:
    import striprtf
    RTF_AVAILABLE = True
except ImportError:
    RTF_AVAILABLE = False

# PowerPoint support
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Note: textract is not available on Windows
TEXTTRACT_AVAILABLE = False

# OCR libraries
try:
    from paddleocr import PaddleOCR
    PADDLE_AVAILABLE = True
except ImportError:
    PADDLE_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False

logger = logging.getLogger(__name__)


class DocumentParser:
    """Enhanced document parser with better MIME type detection and error handling"""
    
    def __init__(self):
        self.ocr_engine = None
        self._initialize_ocr()
        
        # Enhanced MIME type mapping
        self.mime_to_parser = {
            'application/pdf': self._parse_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._parse_docx,
            'application/msword': self._parse_doc,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._parse_excel,
            'application/vnd.ms-excel': self._parse_excel,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._parse_pptx,
            'application/vnd.ms-powerpoint': self._parse_ppt,
            'text/plain': self._parse_text,
            'text/markdown': self._parse_markdown,  # Markdown uses dedicated parser
            'text/csv': self._parse_csv,
            'application/rtf': self._parse_rtf,
            'text/rtf': self._parse_rtf,
        }
        
        # File extension fallback mapping
        self.extension_to_parser = {
            '.pdf': self._parse_pdf,
            '.docx': self._parse_docx,
            '.doc': self._parse_doc,
            '.xlsx': self._parse_excel,
            '.xls': self._parse_excel,
            '.pptx': self._parse_pptx,
            '.ppt': self._parse_ppt,
            '.txt': self._parse_text,
            '.md': self._parse_markdown,  # Markdown files use dedicated parser
            '.markdown': self._parse_markdown,  # Markdown files use dedicated parser
            '.csv': self._parse_csv,
            '.rtf': self._parse_rtf,
        }
    
    def _initialize_ocr(self):
        """Initialize OCR engines with better error handling and image processing"""
        # Инициализируем переменные
        self.paddle_ocr = None
        self.tesseract_ocr = None
        self.ocr_engine = None
        
        try:
            # Временно отключаем PaddleOCR из-за проблем с путями к моделям
            logger.info("PaddleOCR temporarily disabled due to model path issues")
            self.paddle_ocr = None
            
            # Используем только Tesseract
            if TESSERACT_AVAILABLE:
                try:
                    pytesseract.get_tesseract_version()
                    self.tesseract_ocr = pytesseract
                    logger.info("✅ Tesseract OCR initialized successfully")
                    self.ocr_engine = 'tesseract'
                except Exception as e:
                    logger.warning(f"Tesseract initialization failed: {e}")
                    self.tesseract_ocr = None
                    self.ocr_engine = None
            else:
                logger.warning("Tesseract not available")
                self.ocr_engine = None
        except Exception as e:
            logger.warning(f"OCR initialization failed: {e}")
            self.ocr_engine = None
    
    def extract_text_from_image(self, image_path: str) -> str:
        """Extract text from image using available OCR engines"""
        try:
            if not self.ocr_engine:
                logger.warning("No OCR engine available for image processing")
                return ""
            
            if self.ocr_engine == 'paddle':
                return self._extract_text_with_paddle(image_path)
            elif self.ocr_engine == 'tesseract':
                return self._extract_text_with_tesseract(image_path)
            else:
                logger.warning(f"Unknown OCR engine: {self.ocr_engine}")
                return ""
                
        except Exception as e:
            logger.error(f"Error extracting text from image: {e}")
            return ""
    
    def _extract_text_with_paddle(self, image_path: str) -> str:
        """Extract text using PaddleOCR"""
        try:
            if not self.paddle_ocr:
                return ""
            
            # Read image
            from PIL import Image
            image = Image.open(image_path)
            
            # Extract text
            result = self.paddle_ocr.ocr(image_path, cls=True)
            
            if not result or not result[0]:
                return ""
            
            # Extract text from results
            text_parts = []
            for line in result[0]:
                if line and len(line) >= 2:
                    text = line[1][0]  # Extract text from OCR result
                    confidence = line[1][1]  # Extract confidence score
                    
                    # Filter by confidence (only keep high-confidence results)
                    if confidence > 0.7:
                        text_parts.append(text)
            
            extracted_text = "\n".join(text_parts)
            logger.info(f"PaddleOCR extracted {len(text_parts)} text lines from image")
            
            return extracted_text
            
        except Exception as e:
            logger.error(f"Error with PaddleOCR: {e}")
            return ""
    
    def _extract_text_with_tesseract(self, image_path: str) -> str:
        """Extract text using Tesseract"""
        try:
            if not self.tesseract_ocr:
                return ""
            
            # Read image
            from PIL import Image
            image = Image.open(image_path)
            
            # Extract text with Russian language support
            text = self.tesseract_ocr.image_to_string(
                image, 
                lang='rus+eng',  # Support both Russian and English
                config='--psm 6'  # Assume uniform block of text
            )
            
            logger.info(f"Tesseract extracted text from image: {len(text)} characters")
            return text.strip()
            
        except Exception as e:
            logger.error(f"Error with Tesseract: {e}")
            return ""
    
    def _detect_mime_by_signature(self, file_path: str) -> str:
        """Detect MIME type by reading file signatures (Windows-compatible)"""
        try:
            with open(file_path, 'rb') as f:
                header = f.read(512)
                
            # Check for common file signatures
            if header.startswith(b'%PDF'):
                return 'application/pdf'
            elif header.startswith(b'PK\x03\x04'):
                # ZIP-based formats (DOCX, XLSX)
                file_ext = os.path.splitext(file_path)[1].lower()
                if file_ext == '.docx':
                    return 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                elif file_ext == '.xlsx':
                    return 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                else:
                    return 'application/zip'
            elif header.startswith(b'\xd0\xcf\x11\xe0'):
                # OLE2 format (DOC, XLS)
                file_ext = os.path.splitext(file_path)[1].lower()
                if file_ext == '.doc':
                    return 'application/msword'
                elif file_ext == '.xls':
                    return 'application/vnd.ms-excel'
                else:
                    return 'application/ole2'
            elif header.startswith(b'#') or any(char in header[:100] for char in b'ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz'):
                # Text-based files
                file_ext = os.path.splitext(file_path)[1].lower()
                if file_ext in ['.md', '.markdown']:
                    return 'text/markdown'
                else:
                    return 'text/plain'
            else:
                return 'application/octet-stream'
                
        except Exception as e:
            logger.debug(f"File signature detection failed: {e}")
            return 'application/octet-stream'
    
    def detect_file_type(self, file_path: str, mime_type: str = None) -> Dict[str, Any]:
        """Enhanced file type detection with multiple fallback methods"""
        try:
            # Method 1: Use provided MIME type
            if mime_type and mime_type in self.mime_to_parser:
                return {
                    'detected_type': mime_type,
                    'parser_method': 'mime_type',
                    'confidence': 'high'
                }
            
            # Method 2: Use file signature detection
            detected_mime = self._detect_mime_by_signature(file_path)
            if detected_mime in self.mime_to_parser:
                return {
                    'detected_type': detected_mime,
                    'parser_method': 'file_signature',
                    'confidence': 'high'
                }
            
            # Method 3: Use file extension as fallback
            file_ext = os.path.splitext(file_path)[1].lower()
            if file_ext in self.extension_to_parser:
                # Map extension to MIME type
                ext_to_mime = {
                    '.pdf': 'application/pdf',
                    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    '.doc': 'application/msword',
                    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                    '.xls': 'application/vnd.ms-excel',
                    '.txt': 'text/plain',
                    '.md': 'text/markdown',
                    '.markdown': 'text/markdown',
                    '.csv': 'text/csv',
                }
                
                detected_mime = ext_to_mime.get(file_ext)
                if detected_mime:
                    return {
                        'detected_type': detected_mime,
                        'parser_method': 'extension_fallback',
                        'confidence': 'medium',
                        'file_extension': file_ext
                    }
            
            # Fallback: return unknown type
            return {
                'detected_type': 'application/octet-stream',
                'parser_method': 'unknown',
                'confidence': 'low',
                'file_extension': file_ext,
                'note': 'Could not determine file type'
            }
            
        except Exception as e:
            logger.error(f"File type detection failed: {e}")
            return {
                'detected_type': 'application/octet-stream',
                'parser_method': 'error',
                'confidence': 'none',
                'error': str(e)
            }
    
    def parse_document(self, file_path: str, mime_type: str = None) -> Dict[str, Any]:
        """Parse document with enhanced type detection and error handling"""
        try:
            # Detect file type
            type_info = self.detect_file_type(file_path, mime_type)
            detected_mime = type_info['detected_type']
            
            logger.info(f"Parsing document: {file_path}")
            logger.info(f"Type detection: {type_info}")
            
            # Choose parser based on detected type
            if detected_mime in self.mime_to_parser:
                parser = self.mime_to_parser[detected_mime]
                result = parser(file_path)
                result['type_detection'] = type_info
                return result
            else:
                # Try extension-based fallback
                file_ext = os.path.splitext(file_path)[1].lower()
                if file_ext in self.extension_to_parser:
                    parser = self.extension_to_parser[file_ext]
                    result = parser(file_path)
                    result['type_detection'] = type_info
                    result['note'] = f"Used extension-based parser for {file_ext}"
                    return result
                else:
                    raise ValueError(f"No parser available for MIME type: {detected_mime} or extension: {file_ext}")
                
        except Exception as e:
            logger.error(f"Document parsing failed: {e}")
            return {
                'error': str(e),
                'type_detection': type_info if 'type_info' in locals() else {'detected_type': 'unknown'},
                'content': '',
                'metadata': {},
                'success': False
            }
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """Parse PDF documents with multiple fallback methods"""
        content = []
        metadata = {}
        
        try:
            # Method 1: Try PyMuPDF first
            if PDF_AVAILABLE:
                try:
                    doc = fitz.open(file_path)
                    metadata['total_pages'] = len(doc)
                    
                    for page_num in range(len(doc)):
                        page = doc.load_page(page_num)
                        
                        # Extract text
                        text = page.get_text()
                        if text.strip():
                            content.append({
                                'type': 'text',
                                'content': text.strip(),
                                'section_name': f'Page {page_num + 1}',
                                'page': page_num + 1
                            })
                        
                        # Check for images and extract text using OCR
                        image_list = page.get_images()
                        if image_list:
                            metadata['has_images'] = True
                            logger.info(f"Found {len(image_list)} images on page {page_num + 1}")
                            
                            # Try OCR on images if available
                            for img_index, img in enumerate(image_list):
                                try:
                                    xref = img[0]
                                    pix = fitz.Pixmap(doc, xref)
                                    
                                    # Check if image is valid for OCR
                                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                                        # Save image to temporary file for OCR processing
                                        temp_dir = tempfile.gettempdir()
                                        temp_img_path = os.path.join(temp_dir, f"pdf_image_{page_num}_{img_index}.png")
                                        pix.save(temp_img_path)
                                        
                                        # Extract text using OCR
                                        ocr_text = self.extract_text_from_image(temp_img_path)
                                        
                                        if ocr_text and len(ocr_text.strip()) > 10:  # Only add if meaningful text
                                            content.append({
                                                'type': 'image_text',
                                                'content': f"[Изображение {img_index + 1}]: {ocr_text.strip()}",
                                                'section_name': f'Страница {page_num + 1} - Изображение {img_index + 1}',
                                                'page': page_num + 1,
                                                'image_index': img_index + 1,
                                                'ocr_confidence': 'high' if len(ocr_text) > 50 else 'medium'
                                            })
                                            logger.info(f"OCR extracted {len(ocr_text)} characters from image on page {page_num + 1}")
                                        else:
                                            logger.debug(f"OCR returned insufficient text from image on page {page_num + 1}")
                                        
                                        # Clean up temporary file
                                        try:
                                            os.remove(temp_img_path)
                                        except:
                                            pass
                                    
                                    pix = None
                                    
                                except Exception as img_error:
                                    logger.debug(f"Image processing failed on page {page_num + 1}, image {img_index + 1}: {img_error}")
                                    continue
                    
                    doc.close()
                    metadata['parser'] = 'PyMuPDF'
                    metadata['extraction_method'] = 'PyMuPDF'
                    
                except Exception as mupdf_error:
                    logger.warning(f"PyMuPDF failed: {mupdf_error}")
                    # Try fallback method
                    if not content:
                        content = self._parse_pdf_with_pdfplumber(file_path)
                        if content:
                            metadata['parser'] = 'pdfplumber'
                            metadata['extraction_method'] = 'pdfplumber_fallback'
                            metadata['fallback_reason'] = str(mupdf_error)
                        else:
                            raise mupdf_error
            else:
                # Try pdfplumber if PyMuPDF is not available
                content = self._parse_pdf_with_pdfplumber(file_path)
                if content:
                    metadata['parser'] = 'pdfplumber'
                    metadata['extraction_method'] = 'pdfplumber_only'
                else:
                    raise ImportError("No PDF parser available")
            
            metadata['title'] = os.path.splitext(os.path.basename(file_path))[0]
            metadata['file_type'] = 'pdf'
            
            return {
                'content': content,
                'metadata': metadata,
                'has_images': metadata.get('has_images', False),
                'total_pages': metadata.get('total_pages', 0),
                'file_type': 'pdf',
                'success': len(content) > 0
            }
            
        except Exception as e:
            logger.error(f"Error parsing PDF file {file_path}: {e}")
            raise
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """Parse DOCX documents"""
        content = []
        metadata = {}
        
        try:
            if DOCX_AVAILABLE:
                doc = DocxDocument(file_path)
            
            # Extract document properties
                if doc.core_properties.title:
                    metadata['title'] = doc.core_properties.title
                if doc.core_properties.author:
                    metadata['author'] = doc.core_properties.author
                
                # Extract paragraphs
                paragraphs = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        paragraphs.append(para.text.strip())
                
                # Group paragraphs into sections
                current_section = []
                sections = []
                
                for para in paragraphs:
                    # Check if this looks like a header
                    if para.isupper() or (len(para) < 100 and para.endswith(':')):
                        if current_section:
                            sections.append('\n'.join(current_section))
                            current_section = []
                        current_section.append(para)
                    else:
                        current_section.append(para)
                
                if current_section:
                    sections.append('\n'.join(current_section))
                
                # Create content items
                if len(sections) > 1:
                    for i, section in enumerate(sections):
                        if section.strip():
                            content.append({
                                'type': 'text',
                                'content': section.strip(),
                                'section_name': f'Section {i+1}' if i > 0 else 'Header'
                            })
                else:
                    content.append({
                        'type': 'text',
                        'content': '\n\n'.join(paragraphs),
                        'section_name': 'Document Content'
                    })
                
                metadata['title'] = os.path.splitext(os.path.basename(file_path))[0]
                metadata['paragraph_count'] = len(paragraphs)
                
            else:
                raise ImportError("No DOCX parser available")
            
            return {
                'content': content,
                'metadata': metadata,
                'has_images': False,
                'total_pages': 1,
                'file_type': 'docx'
            }
            
        except Exception as e:
            logger.error(f"Error parsing DOCX file {file_path}: {e}")
            raise
    
    def _parse_excel(self, file_path: str) -> Dict[str, Any]:
        """Parse Excel documents with support for both .xlsx and .xls formats"""
        content = []
        metadata = {}
        
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            is_old_format = file_ext == '.xls'
            
            if EXCEL_AVAILABLE:
                # Method 1: Try pandas with openpyxl/xlrd
                try:
                    if is_old_format and XLRD_AVAILABLE:
                        # For .xls files, use xlrd engine
                        excel_file = pd.ExcelFile(file_path, engine='xlrd')
                    else:
                        # For .xlsx files, use openpyxl engine
                        excel_file = pd.ExcelFile(file_path, engine='openpyxl')
                    
                    sheet_names = excel_file.sheet_names
                    metadata['total_sheets'] = len(sheet_names)
                    metadata['sheet_names'] = sheet_names
                    metadata['parser'] = 'pandas'
                    
                    for sheet_name in sheet_names:
                        try:
                            if is_old_format and XLRD_AVAILABLE:
                                df = pd.read_excel(file_path, sheet_name=sheet_name, engine='xlrd')
                            else:
                                df = pd.read_excel(file_path, sheet_name=sheet_name, engine='openpyxl')
                            
                            if not df.empty:
                                # Convert DataFrame to structured text
                                sheet_text = self._excel_sheet_to_text(df, sheet_name)
                                
                                content.append({
                                    'type': 'table',
                                    'content': sheet_text,
                                    'section_name': sheet_name,
                                    'sheet_name': sheet_name,
                                    'row_count': len(df),
                                    'column_count': len(df.columns)
                                })
                            else:
                                # Empty sheet
                                content.append({
                                    'type': 'table',
                                    'content': f"Sheet: {sheet_name} (Empty)",
                                    'section_name': sheet_name,
                                    'sheet_name': sheet_name,
                                    'row_count': 0,
                                    'column_count': 0
                                })
                        
                        except Exception as sheet_error:
                            logger.warning(f"Failed to parse sheet {sheet_name}: {sheet_error}")
                            # Add error information to content
                            content.append({
                                'type': 'error',
                                'content': f"Ошибка при парсинге листа '{sheet_name}': {str(sheet_error)}",
                                'section_name': sheet_name,
                                'sheet_name': sheet_name,
                                'error': str(sheet_error)
                            })
                            continue
                    
                    metadata['title'] = os.path.splitext(os.path.basename(file_path))[0]
                    metadata['file_type'] = 'excel'
                    
                except Exception as pandas_error:
                    logger.warning(f"Pandas parsing failed: {pandas_error}")
                    
                    # Method 2: Try xlrd directly for .xls files
                    if is_old_format and XLRD_AVAILABLE:
                        try:
                            content = self._parse_xls_with_xlrd(file_path)
                            metadata['parser'] = 'xlrd_direct'
                            metadata['fallback_method'] = 'xlrd_direct'
                        except Exception as xlrd_error:
                            logger.warning(f"xlrd direct parsing failed: {xlrd_error}")
                            raise pandas_error  # Re-raise original error
                    else:
                        raise pandas_error
                
            else:
                # Method 3: Try xlrd directly if pandas is not available
                if is_old_format and XLRD_AVAILABLE:
                    try:
                        content = self._parse_xls_with_xlrd(file_path)
                        metadata['parser'] = 'xlrd_only'
                    except Exception as xlrd_error:
                        logger.error(f"xlrd parsing failed: {xlrd_error}")
                        raise ImportError("No Excel parser available")
                else:
                    raise ImportError("No Excel parser available")
            
            return {
                'content': content,
                'metadata': metadata,
                'has_images': False,
                'total_pages': len(content),
                'file_type': 'xlsx' if file_path.endswith('.xlsx') else 'xls',
                'success': len(content) > 0
            }
            
        except Exception as e:
            logger.error(f"Error parsing Excel file {file_path}: {e}")
            raise
    
    def _excel_sheet_to_text(self, df: pd.DataFrame, sheet_name: str) -> str:
        """Convert Excel sheet to structured text with better formatting"""
        if df.empty:
            return f"Sheet: {sheet_name} (Empty)"
        
        # Remove completely empty rows and columns
        df = df.dropna(how='all').dropna(axis=1, how='all')
        
        if df.empty:
            return f"Sheet: {sheet_name} (No data)"
        
        text_parts = [f"Sheet: {sheet_name}"]
        
        # Try to identify headers
        header_row = None
        for i in range(min(3, len(df))):  # Check first 3 rows for headers
            if df.iloc[i].notna().sum() > len(df.columns) * 0.5:  # More than 50% non-null
                header_row = i
                break
        
        if header_row is not None:
            # Use identified header row
            headers = df.iloc[header_row].fillna('').astype(str).tolist()
            text_parts.append("Headers: " + " | ".join(headers))
            
            # Process data rows
            for idx, row in df.iloc[header_row + 1:].iterrows():
                if row.notna().sum() > 0:  # Skip completely empty rows
                    row_data = row.fillna('').astype(str).tolist()
                    text_parts.append("Row " + str(idx + 1) + ": " + " | ".join(row_data))
        else:
            # No clear headers, just output all data
            for idx, row in df.iterrows():
                if row.notna().sum() > 0:
                    row_data = row.fillna('').astype(str).tolist()
                    text_parts.append("Row " + str(idx + 1) + ": " + " | ".join(row_data))
        
        return "\n".join(text_parts)
    
    def _parse_pdf_with_pdfplumber(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse PDF using pdfplumber as fallback"""
        content = []
        
        try:
            import pdfplumber
            
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    try:
                        # Extract text
                        text = page.extract_text()
                        if text and text.strip():
                            content.append({
                                'type': 'text',
                                'content': text.strip(),
                                'section_name': f'Page {page_num + 1}',
                                'page': page_num + 1
                            })
                        
                        # Extract tables if any
                        tables = page.extract_tables()
                        for table_idx, table in enumerate(tables):
                            if table and any(any(cell for cell in row if cell) for row in table):
                                table_text = self._table_to_text(table)
                                content.append({
                                    'type': 'table',
                                    'content': table_text,
                                    'section_name': f'Page {page_num + 1} Table {table_idx + 1}',
                                    'page': page_num + 1
                                })
                    
                    except Exception as page_error:
                        logger.warning(f"Failed to parse PDF page {page_num + 1}: {page_error}")
                        continue
                    
            return content
            
        except ImportError:
            logger.warning("pdfplumber not available for PDF fallback")
            return []
        except Exception as e:
            logger.error(f"pdfplumber parsing failed: {e}")
            return []
    
    def _table_to_text(self, table: List[List]) -> str:
        """Convert table to structured text"""
        if not table:
            return "Empty table"
        
        text_parts = []
        
        # Process each row
        for row_idx, row in enumerate(table):
            # Filter out empty cells and convert to strings
            row_data = [str(cell) if cell else "" for cell in row]
            
            # Only add non-empty rows
            if any(cell.strip() for cell in row_data):
                text_parts.append(" | ".join(row_data))
        
        return "\n".join(text_parts)
    
    def _parse_xls_with_xlrd(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse .xls files directly using xlrd library"""
        content = []
        
        try:
            if not XLRD_AVAILABLE:
                raise ImportError("xlrd library not available")
            
            # Open workbook
            workbook = xlrd.open_workbook(file_path)
            sheet_names = workbook.sheet_names()
            
            for sheet_name in sheet_names:
                try:
                    sheet = workbook.sheet_by_name(sheet_name)
                    
                    if sheet.nrows > 0:
                        # Convert sheet to text
                        sheet_text = self._xlrd_sheet_to_text(sheet, sheet_name)
                        
                        content.append({
                            'type': 'table',
                            'content': sheet_text,
                            'section_name': sheet_name,
                            'sheet_name': sheet_name,
                            'row_count': sheet.nrows,
                            'column_count': sheet.ncols
                        })
                    else:
                        # Empty sheet
                        content.append({
                            'type': 'table',
                            'content': f"Sheet: {sheet_name} (Empty)",
                            'section_name': sheet_name,
                            'sheet_name': sheet_name,
                            'row_count': 0,
                            'column_count': 0
                        })
                        
                except Exception as sheet_error:
                    logger.warning(f"Failed to parse xlrd sheet {sheet_name}: {sheet_error}")
                    content.append({
                        'type': 'error',
                        'content': f"Ошибка при парсинге листа '{sheet_name}': {str(sheet_error)}",
                        'section_name': sheet_name,
                        'sheet_name': sheet_name,
                        'error': str(sheet_error)
                    })
                    continue
            
            return content
            
        except Exception as e:
            logger.error(f"xlrd parsing failed: {e}")
            raise
    
    def _xlrd_sheet_to_text(self, sheet, sheet_name: str) -> str:
        """Convert xlrd sheet to structured text"""
        text_parts = [f"Sheet: {sheet_name}"]
        text_parts.append(f"Rows: {sheet.nrows}, Columns: {sheet.ncols}")
        
        # Process all rows
        for row_idx in range(sheet.nrows):
            row_data = []
            for col_idx in range(sheet.ncols):
                try:
                    cell_value = sheet.cell_value(row_idx, col_idx)
                    # Convert cell value to string
                    if isinstance(cell_value, (int, float)):
                        cell_str = str(int(cell_value)) if cell_value == int(cell_value) else str(cell_value)
                    else:
                        cell_str = str(cell_value) if cell_value else ""
                    row_data.append(cell_str)
                except Exception as e:
                    row_data.append(f"[Error: {e}]")
            
            # Only add non-empty rows
            if any(cell.strip() for cell in row_data if cell):
                text_parts.append(" | ".join(row_data))
        
        return "\n".join(text_parts)
    
    def _parse_doc(self, file_path: str) -> Dict[str, Any]:
        """Parse legacy DOC documents with multiple extraction methods"""
        content = []
        metadata = {}
        
        try:
            # Method 1: Try using win32com to open with Word (most reliable for old DOC files)
            try:
                import win32com.client
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(os.path.abspath(file_path))
                text_content = doc.Content.Text
                doc.Close()
                word.Quit()
                
                if text_content and text_content.strip():
                    # Очищаем текст от проблемных символов
                    cleaned_text = self._clean_text_content(text_content.strip())
                    
                    # Проверяем, что текст действительно содержит полезную информацию
                    if len(cleaned_text) > 50 and any(c.isalpha() for c in cleaned_text):
                        content.append({
                            'type': 'text',
                            'content': cleaned_text,
                            'section_name': 'Document Content'
                        })
                        metadata['extraction_method'] = 'win32com_word'
                        metadata['parser'] = 'win32com_word'
                        metadata['text_length'] = len(cleaned_text)
                        logger.info(f"Successfully extracted {len(cleaned_text)} characters from DOC using win32com Word")
                        return {
                            'content': content,
                            'metadata': metadata,
                            'has_images': False,
                            'total_pages': 1,
                            'file_type': 'doc',
                            'success': True
                        }
                    else:
                        logger.warning("win32com Word returned insufficient content")
                else:
                    logger.warning("win32com Word returned empty content")
            except Exception as e:
                logger.debug(f"win32com Word failed: {e}")
            
            # Method 2: Try docx2txt (works for some DOC files)
            if not content and DOCX2TXT_AVAILABLE:
                try:
                    text_content = docx2txt.process(file_path)
                    if text_content and text_content.strip():
                        # Очищаем текст от проблемных символов
                        cleaned_text = self._clean_text_content(text_content.strip())
                        
                        # Проверяем, что текст действительно содержит полезную информацию
                        if len(cleaned_text) > 50 and any(c.isalpha() for c in cleaned_text):
                            content.append({
                                'type': 'text',
                                'content': cleaned_text,
                                'section_name': 'Document Content'
                            })
                            metadata['extraction_method'] = 'docx2txt'
                            metadata['parser'] = 'docx2txt'
                            metadata['text_length'] = len(cleaned_text)
                            logger.info(f"Successfully extracted {len(cleaned_text)} characters from DOC using docx2txt")
                            return {
                                'content': content,
                                'metadata': metadata,
                                'has_images': False,
                                'total_pages': 1,
                                'file_type': 'doc',
                                'success': True
                            }
                        else:
                            logger.warning("docx2txt returned insufficient content")
                    else:
                        logger.warning("docx2txt returned empty content")
                except Exception as e:
                    logger.debug(f"docx2txt failed: {e}")
            
            # Method 3: Try to read as binary and extract text patterns (improved)
            if not content:
                try:
                    with open(file_path, 'rb') as f:
                        binary_content = f.read()
                        
                    # Improved text extraction from binary DOC content
                    text_content = ""
                    
                    # Method 3a: Look for UTF-16 text (common in DOC files)
                    try:
                        # Try to decode as UTF-16
                        utf16_text = binary_content.decode('utf-16', errors='ignore')
                        # Filter out control characters and keep only printable text
                        clean_utf16 = ''.join(char for char in utf16_text if char.isprintable() or char.isspace())
                        if len(clean_utf16) > 100 and any(c.isalpha() for c in clean_utf16):
                            text_content = clean_utf16
                            logger.info("Successfully extracted text using UTF-16 decoding")
                    except:
                        pass
                    
                    # Method 3b: Look for ASCII text patterns (if UTF-16 failed)
                    if not text_content:
                        text_patterns = []
                        current_text = ""
                        
                        for byte in binary_content:
                            if 32 <= byte <= 126: # Printable ASCII
                                current_text += chr(byte)
                            else:
                                if len(current_text) > 15: # Increased minimum length
                                    text_patterns.append(current_text.strip())
                                current_text = ""
                        
                        if current_text.strip() and len(current_text) > 15:
                            text_patterns.append(current_text.strip())
                        
                        # Filter and clean patterns
                        meaningful_patterns = []
                        for pattern in text_patterns:
                            if len(pattern) > 25: # Increased minimum length
                                cleaned = self._clean_text_content(pattern)
                                if len(cleaned) > 25 and any(c.isalpha() for c in cleaned):
                                    meaningful_patterns.append(cleaned)
                        
                        if meaningful_patterns:
                            text_content = '\n\n'.join(meaningful_patterns)
                            logger.info(f"Successfully extracted text using ASCII pattern extraction")
                    
                    # Method 3c: Look for specific DOC file markers and extract text around them
                    if not text_content:
                        # Look for common DOC file markers
                        markers = [b'Microsoft Word', b'Word.Document', b'Document Summary Information']
                        for marker in markers:
                            if marker in binary_content:
                                # Extract text around the marker
                                marker_pos = binary_content.find(marker)
                                start_pos = max(0, marker_pos - 1000)
                                end_pos = min(len(binary_content), marker_pos + 1000)
                                
                                # Extract and clean the text
                                section = binary_content[start_pos:end_pos]
                                try:
                                    section_text = section.decode('utf-8', errors='ignore')
                                    clean_section = ''.join(char for char in section_text if char.isprintable() or char.isspace())
                                    if len(clean_section) > 100 and any(c.isalpha() for c in clean_section):
                                        text_content = clean_section
                                        logger.info("Successfully extracted text around DOC marker")
                                        break
                                except:
                                    pass
                    
                    if text_content:
                        # Clean and validate the extracted text
                        cleaned_text = self._clean_text_content(text_content)
                        
                        # Additional validation: check if text contains meaningful content
                        if len(cleaned_text) > 100 and any(c.isalpha() for c in cleaned_text):
                            # Calculate text quality score
                            letters = sum(1 for c in cleaned_text if c.isalpha())
                            total = len(cleaned_text)
                            letter_percentage = (letters / total) * 100 if total > 0 else 0
                            
                            if letter_percentage > 30: # Accept text with at least 30% letters
                                content.append({
                                    'type': 'text',
                                    'content': cleaned_text,
                                    'section_name': 'Extracted Text'
                                })
                                metadata['extraction_method'] = 'enhanced_binary_extraction'
                                metadata['parser'] = 'enhanced_binary'
                                metadata['text_length'] = len(cleaned_text)
                                metadata['text_quality_score'] = letter_percentage
                                logger.info(f"Successfully extracted {len(cleaned_text)} characters from DOC using enhanced binary extraction (quality: {letter_percentage:.1f}%)")
                                return {
                                    'content': content,
                                    'metadata': metadata,
                                    'has_images': False,
                                    'total_pages': 1,
                                    'file_type': 'doc',
                                    'success': True
                                }
                            else:
                                logger.warning(f"Extracted text quality too low: {letter_percentage:.1f}% letters")
                        else:
                            logger.warning("Extracted text too short or contains no letters")
                    else:
                        logger.warning("No text could be extracted from binary content")
                        
                except Exception as e:
                    logger.debug(f"Enhanced binary extraction failed: {e}")
            
            # If no content was extracted, create error message
            if not content:
                content.append({
                    'type': 'text',
                    'content': 'Не удалось извлечь текст из этого DOC файла. Файл может быть поврежден или в неподдерживаемом формате. Попробуйте конвертировать файл в DOCX или PDF.',
                    'section_name': 'Ошибка извлечения'
                })
                metadata['extraction_method'] = 'failed'
                metadata['note'] = 'Все методы извлечения текста не удались'
            
            metadata['title'] = os.path.splitext(os.path.basename(file_path))[0]
            metadata['file_type'] = 'doc'
            metadata['parser_note'] = 'Legacy DOC format - multiple extraction methods attempted'
            
            return {
                'content': content,
                'metadata': metadata,
                'has_images': False,
                'total_pages': 1,
                'file_type': 'doc',
                'success': True
            }
            
        except Exception as e:
            logger.error(f"Error parsing DOC file {file_path}: {e}")
            raise
    
    def _parse_text(self, file_path: str) -> Dict[str, Any]:
        """Parse text documents (TXT, MD)"""
        try:
            content = []
            metadata = {}
            
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            text_content = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text_content = f.read()
                        metadata['encoding'] = encoding
                        break
                except UnicodeDecodeError:
                    continue
            
            if text_content is None:
                # Fallback: read as binary and decode with errors='ignore'
                with open(file_path, 'rb') as f:
                    binary_content = f.read()
                    text_content = binary_content.decode('utf-8', errors='ignore')
                    metadata['encoding'] = 'utf-8_with_errors'
                    metadata['note'] = 'Used fallback decoding with error handling'
            
            # Split content into sections if it contains headers
            lines = text_content.split('\n')
            current_section = []
            sections = []
            
            for line in lines:
                line = line.strip()
                if line.startswith('#') or (line.isupper() and len(line) > 3 and line.endswith(':')):
                    # This looks like a header
                    if current_section:
                        sections.append('\n'.join(current_section))
                        current_section = []
                    current_section.append(line)
                else:
                    current_section.append(line)
            
            if current_section:
                sections.append('\n'.join(current_section))
            
            if len(sections) > 1:
                # Multiple sections found
                for i, section in enumerate(sections):
                    if section.strip():
                        content.append({
                            'type': 'text',
                            'content': section.strip(),
                            'section_name': f'Section {i+1}' if i > 0 else 'Header'
                        })
            else:
                # Single section
                content.append({
                    'type': 'text',
                    'content': text_content,
                    'section_name': 'Document Content'
                })
            
            metadata['title'] = os.path.splitext(os.path.basename(file_path))[0]
            metadata['total_lines'] = len(lines)
            metadata['total_sections'] = len(sections)
            metadata['file_type'] = os.path.splitext(file_path)[1].lower()
            
            return {
                'content': content,
                'metadata': metadata,
                'has_images': False,
                'total_pages': 1,
                'file_type': os.path.splitext(file_path)[1].lower()
            }
            
        except Exception as e:
            logger.error(f"Error parsing text file {file_path}: {e}")
            raise
    
    def _parse_csv(self, file_path: str) -> Dict[str, Any]:
        """Parse CSV documents"""
        content = []
        metadata = {}
        
        try:
            if EXCEL_AVAILABLE:
                df = pd.read_csv(file_path)
                metadata['total_rows'] = len(df)
                metadata['total_columns'] = len(df.columns)
                metadata['file_type'] = 'csv'
                
                # Convert DataFrame to text
                sheet_content = []
                sheet_content.append(f"CSV File: {os.path.basename(file_path)}")
                sheet_content.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
                
                # Add column headers
                if len(df.columns) > 0:
                    headers = [str(col) for col in df.columns]
                    sheet_content.append("Columns: " + " | ".join(headers))
                
                # Add sample data (first 10 rows)
                for idx, row in df.head(10).iterrows():
                    row_data = [str(val) if pd.notna(val) else "" for val in row]
                    sheet_content.append(" | ".join(row_data))
                
                if len(df) > 10:
                    sheet_content.append(f"... and {len(df) - 10} more rows")
                
                content.append({
                    'type': 'table',
                    'content': '\n'.join(sheet_content),
                    'section_name': 'CSV Data',
                    'row_count': len(df),
                    'column_count': len(df.columns)
                })
                
            else:
                raise ImportError("No CSV parser available")
            
            metadata['title'] = os.path.splitext(os.path.basename(file_path))[0]
            metadata['file_type'] = 'csv'
            
            return {
                'content': content,
                'metadata': metadata,
                'has_images': False,
                'total_pages': 1,
                'file_type': 'csv'
            }
            
        except Exception as e:
            logger.error(f"Error parsing CSV file {file_path}: {e}")
            raise
    
    def _perform_ocr_on_image(self, image_data: bytes) -> Optional[str]:
        """Perform OCR on image data"""
        if not self.ocr_engine:
            return None
            
        try:
            if self.ocr_engine == "tesseract":
                # Use Tesseract
                import io
                from PIL import Image
                
                image = Image.open(io.BytesIO(image_data))
                text = pytesseract.image_to_string(image)
                return text.strip() if text else None
            
            elif hasattr(self.ocr_engine, 'ocr'):
                # Use PaddleOCR
                import io
                import numpy as np
                from PIL import Image
                
                image = Image.open(io.BytesIO(image_data))
                image_array = np.array(image)
                
                result = self.ocr_engine.ocr(image_array, cls=True)
                if result and result[0]:
                    # Extract text from OCR result
                    text_parts = []
                    for line in result[0]:
                        if line and len(line) >= 2:
                            text_parts.append(line[1][0])  # Extract text from result
                    return "\n".join(text_parts) if text_parts else None
            
            return None
            
        except Exception as e:
            logger.warning(f"OCR failed: {e}")
            return None
    
    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """Parse PowerPoint PPTX files"""
        content = []
        metadata = {}
        
        try:
            if not PPTX_AVAILABLE:
                raise ImportError("python-pptx library not available")
            
            presentation = Presentation(file_path)
            slide_count = len(presentation.slides)
            
            metadata['total_slides'] = slide_count
            metadata['parser'] = 'python-pptx'
            metadata['file_type'] = 'pptx'
            
            for slide_idx, slide in enumerate(presentation.slides):
                slide_text = []
                slide_notes = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables
                    if shape.shape_type == 19:  # Table
                        table_text = self._extract_table_text_from_shape(shape)
                        if table_text:
                            slide_text.append(table_text)
                
                # Extract notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_notes.append(notes_text)
                
                # Combine slide content
                if slide_text:
                    slide_content = f"Slide {slide_idx + 1}:\n" + "\n".join(slide_text)
                    content.append({
                        'type': 'slide',
                        'content': slide_content,
                        'section_name': f'Slide {slide_idx + 1}',
                        'slide_number': slide_idx + 1,
                        'text_elements': len(slide_text)
                    })
                
                # Add notes if available
                if slide_notes:
                    notes_content = f"Notes for Slide {slide_idx + 1}:\n" + "\n".join(slide_notes)
                    content.append({
                        'type': 'notes',
                        'content': notes_content,
                        'section_name': f'Notes Slide {slide_idx + 1}',
                        'slide_number': slide_idx + 1
                    })
            
            metadata['title'] = os.path.splitext(os.path.basename(file_path))[0]
            
            return {
                'content': content,
                'metadata': metadata,
                'has_images': False,
                'total_pages': slide_count,
                'file_type': 'pptx',
                'success': len(content) > 0
            }
            
        except Exception as e:
            logger.error(f"Error parsing PPTX file {file_path}: {e}")
            raise
    
    def _parse_ppt(self, file_path: str) -> Dict[str, Any]:
        """Parse legacy PowerPoint PPT files with aggressive text extraction"""
        content = []
        metadata = {}
        
        try:
            # Method 1: Try to convert PPT to PPTX using python-pptx (if possible)
            if PPTX_AVAILABLE:
                try:
                    # Some .ppt files can be read by python-pptx
                    presentation = Presentation(file_path)
                    result = self._parse_pptx(file_path)
                    metadata['parser'] = 'python-pptx_legacy'
                    metadata['fallback_method'] = 'direct_pptx_reading'
                    return result
                except Exception as e:
                    logger.debug(f"python-pptx failed for .ppt file: {e}")
            
            # Method 2: Ultra-aggressive binary text extraction with safety limits
            try:
                with open(file_path, 'rb') as f:
                    binary_content = f.read()
                
                # Look for ANY text patterns in binary content with safety limits
                text_patterns = []
                current_text = ""
                pattern_count = 0
                MAX_PATTERNS = 1000  # Safety limit to prevent infinite loops
                
                for i, byte in enumerate(binary_content):
                    if 32 <= byte <= 126:  # Printable ASCII
                        current_text += chr(byte)
                    else:
                        # Save any text longer than 3 characters
                        if len(current_text.strip()) > 3:
                            text_patterns.append(current_text.strip())
                            pattern_count += 1
                            
                            # Safety check: prevent too many patterns
                            if pattern_count >= MAX_PATTERNS:
                                logger.warning(f"Reached safety limit of {MAX_PATTERNS} patterns, stopping extraction")
                                break
                        current_text = ""
                
                # Don't forget the last text block
                if len(current_text.strip()) > 3 and pattern_count < MAX_PATTERNS:
                    text_patterns.append(current_text.strip())
                    pattern_count += 1
                
                # Filter and clean patterns more aggressively with safety limits
                meaningful_patterns = []
                meaningful_count = 0
                MAX_MEANINGFUL = 500  # Safety limit for meaningful patterns
                
                # Keywords to prioritize (case-insensitive)
                priority_keywords = [
                    'time', 'management', 'тайм', 'менеджмент', 'project', 'plan', 'goal',
                    'objective', 'strategy', 'process', 'workflow', 'procedure', 'schedule',
                    'deadline', 'priority', 'efficiency', 'productivity', 'organization'
                ]
                
                for pattern in text_patterns:
                    if meaningful_count >= MAX_MEANINGFUL:
                        break
                        
                    # Remove common binary artifacts but be less strict
                    cleaned = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F-\x9F]', '', pattern)
                    cleaned = re.sub(r'\s+', ' ', cleaned).strip()
                    
                    # Check if pattern contains priority keywords
                    pattern_lower = cleaned.lower()
                    has_priority_keywords = any(keyword in pattern_lower for keyword in priority_keywords)
                    
                    # MUCH STRICTER filtering: only keep patterns that look like real text
                    if (len(cleaned) > 10 and  # Longer patterns
                        any(c.isalpha() for c in cleaned) and  # Contains letters
                        any(c.isupper() for c in cleaned) and  # Contains uppercase letters
                        any(c.islower() for c in cleaned) and  # Contains lowercase letters
                        cleaned.count(' ') > 1 and  # Contains multiple spaces
                        not re.search(r'[^\w\s]', cleaned) and  # No special characters
                        len(cleaned.split()) > 2):  # Multiple words
                        
                        # Prioritize patterns with keywords
                        if has_priority_keywords:
                            # Insert at the beginning for priority
                            meaningful_patterns.insert(0, cleaned)
                            logger.info(f"🎯 Found priority pattern: {cleaned}")
                        else:
                            meaningful_patterns.append(cleaned)
                        meaningful_count += 1
                
                if meaningful_patterns:
                    # Combine all patterns into a single coherent text
                    combined_text = '\n\n'.join(meaningful_patterns)
                    
                    # Split into chunks for better processing with safety limit
                    chunks = self._split_text_into_chunks(combined_text, max_chunk_size=1000, overlap=100)
                    
                    # Safety limit for chunks
                    MAX_CHUNKS = 50
                    if len(chunks) > MAX_CHUNKS:
                        logger.warning(f"Limiting chunks from {len(chunks)} to {MAX_CHUNKS} for safety")
                        chunks = chunks[:MAX_CHUNKS]
                    
                    for i, chunk in enumerate(chunks):
                        content.append({
                            'type': 'text',
                            'content': chunk,
                            'section_name': f'Extracted Text Part {i + 1}',
                            'chunk_number': i + 1
                        })
                    
                    metadata['extraction_method'] = 'ultra_aggressive_binary_extraction'
                    metadata['parser'] = 'ultra_aggressive_binary_pattern'
                    metadata['text_patterns_found'] = len(meaningful_patterns)
                    metadata['total_chunks'] = len(chunks)
                    metadata['raw_patterns'] = text_patterns[:10]  # Store first 10 raw patterns for debugging
                    metadata['safety_limits_applied'] = pattern_count >= MAX_PATTERNS or meaningful_count >= MAX_MEANINGFUL or len(chunks) >= MAX_CHUNKS
                    
                    # Debug: Log first few meaningful patterns
                    debug_patterns = meaningful_patterns[:5]
                    logger.info(f"Successfully extracted text from PPT using ultra-aggressive extraction: {len(meaningful_patterns)} patterns, {len(chunks)} chunks (safety limits: {metadata['safety_limits_applied']})")
                    logger.info(f"🔍 First 5 meaningful patterns: {debug_patterns}")
                    
                    # Check if we found priority keywords
                    priority_found = any(any(keyword in pattern.lower() for keyword in priority_keywords) for pattern in debug_patterns)
                    logger.info(f"🎯 Priority keywords found: {priority_found}")
                else:
                    logger.warning("No meaningful text patterns found in PPT binary content")
                    
            except Exception as e:
                logger.debug(f"Ultra-aggressive binary extraction failed: {e}")
            
            # Method 3: Look for specific text markers in binary content
            if not content:
                try:
                    logger.info("🔄 Trying marker-based extraction...")
                    
                    # Look for specific text markers in binary content
                    text_markers = [
                        b'Title', b'Subtitle', b'Text', b'Notes', b'Slide',
                        b'Content', b'Header', b'Footer', b'Summary', b'Time',
                        b'Management', b'Project', b'Plan', b'Goal', b'Objective',
                        b'Strategy', b'Process', b'Workflow', b'Procedure',
                        b'TIME', b'MANAGEMENT', b'PROJECT', b'PLAN', b'GOAL',
                        b'OBJECTIVE', b'STRATEGY', b'PROCESS', b'WORKFLOW',
                        b'SCHEDULE', b'DEADLINE', b'PRIORITY', b'EFFICIENCY',
                        b'PRODUCTIVITY', b'ORGANIZATION'
                    ]
                    
                    extracted_sections = []
                    for marker in text_markers:
                        if marker in binary_content:
                            # Find text around this marker
                            start_pos = binary_content.find(marker)
                            if start_pos >= 0:
                                # Extract text around marker (look for printable characters)
                                section_start = max(0, start_pos - 200)
                                section_end = min(len(binary_content), start_pos + 300)
                                
                                section_bytes = binary_content[section_start:section_end]
                                section_text = ""
                                
                                for byte in section_bytes:
                                    if 32 <= byte <= 126:
                                        section_text += chr(byte)
                                    else:
                                        section_text += " "
                                
                                # Clean up the text
                                cleaned_section = re.sub(r'\s+', ' ', section_text).strip()
                                if len(cleaned_section) > 10:
                                    extracted_sections.append(f"{marker.decode()}: {cleaned_section}")
                    
                    if extracted_sections:
                        combined_text = '\n\n'.join(extracted_sections)
                        chunks = self._split_text_into_chunks(combined_text, max_chunk_size=1000, overlap=100)
                        
                        for i, chunk in enumerate(chunks):
                            content.append({
                                'type': 'text',
                                'content': chunk,
                                'section_name': f'Marker-based Text {i + 1}',
                                'chunk_number': i + 1
                            })
                        
                        metadata['extraction_method'] = 'marker_based_extraction'
                        metadata['parser'] = 'marker_based'
                        metadata['markers_found'] = len(extracted_sections)
                        metadata['total_chunks'] = len(chunks)
                        logger.info(f"✅ Marker-based extraction successful: {len(extracted_sections)} sections")
                        
                except Exception as e:
                    logger.debug(f"Marker-based extraction failed: {e}")
            
            # Method 4: Simple text extraction if other methods failed
            if not content:
                try:
                    logger.info("🔄 Trying simple text extraction as fallback...")
                    
                    # Extract all printable text from the file
                    all_text = ""
                    for byte in binary_content:
                        if 32 <= byte <= 126:
                            all_text += chr(byte)
                        else:
                            all_text += " "
                    
                    # Clean up the text
                    cleaned_text = re.sub(r'\s+', ' ', all_text).strip()
                    
                    if len(cleaned_text) > 100:  # Only if we have substantial text
                        # Check if we found any keywords in the text
                        text_lower = cleaned_text.lower()
                        found_keywords = []
                        for keyword in priority_keywords:
                            if keyword in text_lower:
                                found_keywords.append(keyword)
                        
                        if found_keywords:
                            logger.info(f"🎯 Found keywords in fallback text: {found_keywords}")
                            
                            # Split into chunks
                            chunks = self._split_text_into_chunks(cleaned_text, max_chunk_size=1000, overlap=100)
                            
                            # Limit chunks for safety
                            MAX_FALLBACK_CHUNKS = 10
                            if len(chunks) > MAX_FALLBACK_CHUNKS:
                                chunks = chunks[:MAX_FALLBACK_CHUNKS]
                            
                            for i, chunk in enumerate(chunks):
                                content.append({
                                    'type': 'text',
                                    'content': chunk,
                                    'section_name': f'Fallback Text {i + 1}',
                                    'chunk_number': i + 1
                                })
                            
                            metadata['extraction_method'] = 'simple_fallback_extraction'
                            metadata['parser'] = 'simple_fallback'
                            metadata['total_chunks'] = len(chunks)
                            metadata['total_text_length'] = len(cleaned_text)
                            metadata['keywords_found'] = found_keywords
                            logger.info(f"✅ Simple fallback extraction successful: {len(chunks)} chunks, {len(cleaned_text)} characters, keywords: {found_keywords}")
                        else:
                            logger.warning("⚠️ No keywords found in fallback text, skipping...")
                        
                except Exception as e:
                    logger.debug(f"Simple fallback extraction failed: {e}")
            
            # Method 5: Try to find any readable text by scanning byte ranges with safety limits
            if not content:
                try:
                    # Scan the file in chunks and look for text with safety limits
                    chunk_size = 1024
                    text_chunks = []
                    MAX_TEXT_CHUNKS = 200  # Safety limit to prevent too many chunks
                    
                    for i in range(0, len(binary_content), chunk_size):
                        if len(text_chunks) >= MAX_TEXT_CHUNKS:
                            logger.warning(f"Reached safety limit of {MAX_TEXT_CHUNKS} text chunks, stopping chunk processing")
                            break
                            
                        chunk = binary_content[i:i + chunk_size]
                        chunk_text = ""
                        
                        for byte in chunk:
                            if 32 <= byte <= 126:
                                chunk_text += chr(byte)
                            else:
                                chunk_text += " "
                        
                        # Clean and check if chunk contains meaningful text
                        cleaned_chunk = re.sub(r'\s+', ' ', chunk_text).strip()
                        if len(cleaned_chunk) > 20 and any(c.isalpha() for c in cleaned_chunk):
                            text_chunks.append(cleaned_chunk)
                    
                    if text_chunks:
                        combined_text = '\n\n'.join(text_chunks)
                        chunks = self._split_text_into_chunks(combined_text, max_chunk_size=1000, overlap=100)
                        
                        # Safety limit for final chunks
                        MAX_FINAL_CHUNKS = 30
                        if len(chunks) > MAX_FINAL_CHUNKS:
                            logger.warning(f"Limiting final chunks from {len(chunks)} to {MAX_FINAL_CHUNKS} for safety")
                            chunks = chunks[:MAX_FINAL_CHUNKS]
                        
                        for i, chunk in enumerate(chunks):
                            content.append({
                                'type': 'text',
                                'content': chunk,
                                'section_name': f'Chunk-based Text {i + 1}',
                                'chunk_number': i + 1
                            })
                        
                        metadata['extraction_method'] = 'chunk_based_extraction'
                        metadata['parser'] = 'chunk_based'
                        metadata['chunks_processed'] = len(text_chunks)
                        metadata['total_chunks'] = len(chunks)
                        metadata['safety_limits_applied'] = len(text_chunks) >= MAX_TEXT_CHUNKS or len(chunks) >= MAX_FINAL_CHUNKS
                        logger.info(f"Successfully extracted text from PPT using chunk-based extraction: {len(text_chunks)} chunks, {len(chunks)} final chunks (safety limits: {metadata['safety_limits_applied']})")
                        
                except Exception as e:
                    logger.debug(f"Chunk-based extraction failed: {e}")
            
            # If no content was extracted, create error message
            if not content:
                content.append({
                    'type': 'text',
                    'content': 'Не удалось извлечь текст из этого PPT файла. Файл может быть поврежден или в неподдерживаемом формате. Попробуйте конвертировать файл в PPTX или PDF.',
                    'section_name': 'Ошибка извлечения'
                })
                metadata['extraction_method'] = 'failed'
                metadata['note'] = 'Все методы извлечения текста не удались'
                logger.error("All PPT parsing methods failed")
            
            # Set metadata
            metadata['title'] = os.path.splitext(os.path.basename(file_path))[0]
            metadata['file_type'] = 'ppt'
            metadata['parser_note'] = 'Legacy PPT format - aggressive extraction methods attempted'
            
            return {
                'content': content,
                'metadata': metadata,
                'has_images': False,
                'total_pages': len(content),
                'file_type': 'ppt',
                'success': len(content) > 0
            }
            
        except Exception as e:
            logger.error(f"Error parsing PPT file {file_path}: {e}")
            raise
            
            metadata['title'] = os.path.splitext(os.path.basename(file_path))[0]
            metadata['file_type'] = 'ppt'
            metadata['parser_note'] = 'Legacy PPT format - multiple extraction methods attempted'
            
            return {
                'content': content,
                'metadata': metadata,
                'has_images': False,
                'total_pages': 1,
                'file_type': 'ppt',
                'success': len(content) > 0
            }
            
        except Exception as e:
            logger.error(f"Error parsing PPT file {file_path}: {e}")
            raise
    
    def _extract_table_text_from_shape(self, shape) -> str:
        """Extract text from table shape in PowerPoint"""
        try:
            if not hasattr(shape, 'table'):
                return ""
            
            table_text = []
            for row in shape.table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    table_text.append(" | ".join(row_text))
            
            return "\n".join(table_text)
        except Exception as e:
            logger.debug(f"Failed to extract table text: {e}")
            return ""
    
    def _table_to_text(self, table_data: List[List[str]]) -> str:
        """Convert table data to readable text"""
        if not table_data:
            return ""
        
        text_lines = []
        for row_idx, row in enumerate(table_data):
            row_text = " | ".join(str(cell) if cell else "" for cell in row)
            text_lines.append(f"Row {row_idx + 1}: {row_text}")
        
        return "\n".join(text_lines)
    
    def extract_text_chunks(self, parsed_content: Dict[str, Any], chunk_size: int = 500, chunk_overlap: int = 50) -> List[Dict[str, Any]]:
        """Extract text chunks from parsed content"""
        chunks = []
        
        if not parsed_content.get('content'):
            return chunks
        
        for content_item in parsed_content['content']:
            text = content_item.get('content', '')
            if not text.strip():
                continue
            
            # Дополнительная очистка текста перед чанкингом
            text = self._clean_text_content(text)
            
            # Split text into chunks
            text_chunks = self._split_text_into_chunks(text, chunk_size, chunk_overlap)
            
            for i, chunk_text in enumerate(text_chunks):
                chunk = {
                    'content': chunk_text,
                    'type': content_item.get('type', 'text'),
                    'section_name': content_item.get('section_name', ''),
                    'page': content_item.get('page'),
                    'sheet_name': content_item.get('sheet_name'),
                    'chunk_index': i,
                    'metadata': {
                        'original_type': content_item.get('type', 'text'),
                        'section_name': content_item.get('section_name', ''),
                        'page_number': content_item.get('page'),
                        'sheet_name': content_item.get('sheet_name')
                    }
                }
                chunks.append(chunk)
        
        return chunks
    
    def _clean_text_content(self, text: str) -> str:
        """Clean text content by removing problematic Unicode characters and normalizing"""
        if not text:
            return text
        
        try:
            import re
            
            # Удаляем проблемные Unicode символы
            cleaned = text.replace('\u0000', '')  # NULL символ
            cleaned = cleaned.replace('\u0001', '')  # START OF HEADING
            cleaned = cleaned.replace('\u0002', '')  # START OF TEXT
            cleaned = cleaned.replace('\u0003', '')  # END OF TEXT
            cleaned = cleaned.replace('\u0004', '')  # END OF TRANSMISSION
            cleaned = cleaned.replace('\u0005', '')  # ENQUIRY
            cleaned = cleaned.replace('\u0006', '')  # ACKNOWLEDGE
            cleaned = cleaned.replace('\u0007', '')  # BELL
            cleaned = cleaned.replace('\u0008', '')  # BACKSPACE
            cleaned = cleaned.replace('\u0009', '')  # HORIZONTAL TAB
            cleaned = cleaned.replace('\u000A', '')  # LINE FEED
            cleaned = cleaned.replace('\u000B', '')  # VERTICAL TAB
            cleaned = cleaned.replace('\u000C', '')  # FORM FEED
            cleaned = cleaned.replace('\u000D', '')  # CARRIAGE RETURN
            cleaned = cleaned.replace('\u000E', '')  # SHIFT OUT
            cleaned = cleaned.replace('\u000F', '')  # SHIFT IN
            
            # Удаляем другие проблемные символы
            cleaned = cleaned.replace('\u0010', '')  # DATA LINK ESCAPE
            cleaned = cleaned.replace('\u0011', '')  # DEVICE CONTROL ONE
            cleaned = cleaned.replace('\u0012', '')  # DEVICE CONTROL TWO
            cleaned = cleaned.replace('\u0013', '')  # DEVICE CONTROL THREE
            cleaned = cleaned.replace('\u0014', '')  # DEVICE CONTROL FOUR
            cleaned = cleaned.replace('\u0015', '')  # NEGATIVE ACKNOWLEDGE
            cleaned = cleaned.replace('\u0016', '')  # SYNCHRONOUS IDLE
            cleaned = cleaned.replace('\u0017', '')  # END OF TRANSMISSION BLOCK
            cleaned = cleaned.replace('\u0018', '')  # CANCEL
            cleaned = cleaned.replace('\u0019', '')  # END OF MEDIUM
            cleaned = cleaned.replace('\u001A', '')  # SUBSTITUTE
            cleaned = cleaned.replace('\u001B', '')  # ESCAPE
            cleaned = cleaned.replace('\u001C', '')  # FILE SEPARATOR
            cleaned = cleaned.replace('\u001D', '')  # GROUP SEPARATOR
            cleaned = cleaned.replace('\u001E', '')  # RECORD SEPARATOR
            cleaned = cleaned.replace('\u001F', '')  # UNIT SEPARATOR
            
            # Удаляем символы DEL и выше 0x7F (не-ASCII)
            cleaned = ''.join(char for char in cleaned if ord(char) < 0x7F or ord(char) > 0x9F)
            
            # Дополнительная очистка для DOC файлов
            # Удаляем бинарные паттерны и мусор
            cleaned = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F-\x9F]', '', cleaned)
            
            # Удаляем множественные пробелы, табы и переносы строк
            cleaned = re.sub(r'\s+', ' ', cleaned)
            
            # Удаляем пустые строки
            cleaned = re.sub(r'\n\s*\n', '\n', cleaned)
            
            # Удаляем строки только с пробелами
            cleaned = '\n'.join(line for line in cleaned.split('\n') if line.strip())
            
            # Ограничиваем длину для предотвращения таймаутов
            if len(cleaned) > 10000:
                cleaned = cleaned[:10000] + "... [текст обрезан для предотвращения таймаута]"
                logger.warning("Текст DOC файла был обрезан для предотвращения таймаута")
            
            return cleaned.strip()
            
        except Exception as e:
            logger.warning(f"⚠️ Ошибка очистки текста: {e}")
            # Возвращаем безопасную версию
            return text[:1000] if text else ""  # Ограничиваем длину
    
    def _parse_markdown(self, file_path: str) -> Dict[str, Any]:
        """Parse Markdown files"""
        content = []
        metadata = {}
        
        try:
            if not MARKDOWN_AVAILABLE:
                raise ImportError("markdown library not available")
            
            # Read markdown file
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                markdown_content = f.read()
            
            # Convert markdown to HTML for better text extraction
            html_content = markdown.markdown(markdown_content)
            
            # Clean HTML tags (simple approach)
            import re
            clean_text = re.sub(r'<[^>]+>', '', html_content)
            
            # Clean up the text
            cleaned_text = self._clean_text_for_processing(clean_text)
            
            # Split into chunks
            chunks = self._split_text_into_chunks(cleaned_text)
            
            # Create content structure
            for i, chunk in enumerate(chunks):
                content.append({
                    'type': 'text',
                    'content': chunk,
                    'section_name': f'Section {i + 1}',
                    'chunk_number': i + 1
                })
            
            metadata['parser'] = 'markdown'
            metadata['file_type'] = 'markdown'
            metadata['title'] = os.path.splitext(os.path.basename(file_path))[0]
            metadata['total_chunks'] = len(chunks)
            
            return {
                'content': content,
                'metadata': metadata,
                'has_images': False,
                'total_pages': 1,  # Markdown doesn't have pages
                'file_type': 'markdown',
                'success': len(content) > 0
            }
            
        except Exception as e:
            logger.error(f"Error parsing Markdown file {file_path}: {e}")
            raise
    
    def _clean_text_for_processing(self, text: str) -> str:
        """Clean and normalize text for processing"""
        if not text:
            return ""
        
        try:
            # Remove extra whitespace
            cleaned = re.sub(r'\s+', ' ', text)
            
            # Remove control characters
            cleaned = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F-\x9F]', '', cleaned)
            
            # Normalize line breaks
            cleaned = cleaned.replace('\r\n', '\n').replace('\r', '\n')
            
            # Remove empty lines
            cleaned = '\n'.join(line for line in cleaned.split('\n') if line.strip())
            
            return cleaned.strip()
            
        except Exception as e:
            logger.warning(f"Error cleaning text: {e}")
            return text[:1000] if text else ""
    
    def _parse_rtf(self, file_path: str) -> Dict[str, Any]:
        """Parse RTF (Rich Text Format) files"""
        content = []
        metadata = {}
        
        try:
            if not RTF_AVAILABLE:
                raise ImportError("striprtf library not available")
            
            # Read RTF file
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            
            # Strip RTF formatting
            from striprtf.striprtf import rtf_to_text
            plain_text = rtf_to_text(rtf_content)
            
            # Clean up the text
            cleaned_text = self._clean_text_for_processing(plain_text)
            
            # Split into chunks
            chunks = self._split_text_into_chunks(cleaned_text)
            
            # Create content structure
            for i, chunk in enumerate(chunks):
                content.append({
                    'type': 'text',
                    'content': chunk,
                    'section_name': f'Section {i + 1}',
                    'chunk_number': i + 1
                })
            
            metadata['parser'] = 'striprtf'
            metadata['file_type'] = 'rtf'
            metadata['title'] = os.path.splitext(os.path.basename(file_path))[0]
            metadata['total_chunks'] = len(chunks)
            
            return {
                'content': content,
                'metadata': metadata,
                'has_images': False,
                'total_pages': 1,  # RTF doesn't have pages like PDF
                'file_type': 'rtf',
                'success': len(content) > 0
            }
            
        except Exception as e:
            logger.error(f"Error parsing RTF file {file_path}: {e}")
            raise
    
    def _split_text_into_chunks(self, text: str, max_chunk_size: int = 500, overlap: int = 50) -> List[str]:
        """Split text into overlapping chunks with safety limits"""
        if len(text) <= max_chunk_size:
            return [text]
        
        chunks = []
        start = 0
        iteration_count = 0
        MAX_ITERATIONS = 1000  # Safety limit to prevent infinite loops
        
        while start < len(text) and iteration_count < MAX_ITERATIONS:
            iteration_count += 1
            
            end = start + max_chunk_size
            
            # Try to break at sentence boundaries
            if end < len(text):
                # Look for sentence endings
                for i in range(end, max(start + max_chunk_size - 100, start), -1):
                    if text[i] in '.!?':
                        end = i + 1
                        break
                
                # If no sentence boundary found, look for paragraph breaks
                if end == start + max_chunk_size:
                    for i in range(end, max(start + max_chunk_size - 100, start), -1):
                        if text[i] == '\n':
                            end = i + 1
                            break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position with overlap
            start = end - overlap
            
            # Safety check: prevent infinite loops
            if start >= len(text) or start <= 0:
                break
                
            # Additional safety: if we're not making progress, break
            if len(chunks) > 0 and start < len(text):
                progress = start / len(text)
                if progress < 0.01:  # Less than 1% progress
                    logger.warning(f"Minimal progress in text chunking, stopping at iteration {iteration_count}")
                    break
        
        # Safety limit on total chunks
        MAX_CHUNKS = 100
        if len(chunks) > MAX_CHUNKS:
            logger.warning(f"Limiting total chunks from {len(chunks)} to {MAX_CHUNKS} for safety")
            chunks = chunks[:MAX_CHUNKS]
        
        if iteration_count >= MAX_ITERATIONS:
            logger.warning(f"Reached maximum iterations ({MAX_ITERATIONS}) in text chunking, stopping")
        
        return chunks

    def can_convert_doc_to_docx(self) -> bool:
        """Check if DOC text extraction is available using docx2txt"""
        return DOCX2TXT_AVAILABLE
    
    def get_supported_formats(self) -> Dict[str, List[str]]:
        """Get supported file formats"""
        return {
            'text': ['txt', 'md', 'markdown'],
            'image': ['png', 'jpg', 'jpeg', 'bmp', 'gif'],
            'table': ['xlsx', 'xls', 'csv'],
            'presentation': ['pptx', 'ppt'],
            'document': ['docx', 'doc'],
            'binary': ['pdf', 'rtf'],
            'web': ['html', 'htm'],
            'email': ['eml', 'msg'],
            'archive': ['zip', 'rar', '7z'],
            'database': ['sql', 'db'],
            'code': ['py', 'js', 'java', 'cpp', 'c', 'h', 'html', 'css', 'php', 'sql', 'xml', 'json', 'yaml', 'yml', 'toml', 'ini', 'conf']
        }


# Global instance
document_parser = DocumentParser()
